#!/usr/bin/env python3
"""dss-docs — universal document ingestion for DeepSearchStack warehouse.

Extracts text from 9+ document formats, stores as markdown in warehouse.

Usage:
    python3 scripts/dss-docs.py <file> [file...]          local files
    python3 scripts/dss-docs.py --dir ./papers/            all files in directory
    python3 scripts/dss-docs.py --url https://example.com/doc.pdf   remote file

Formats (auto-detected by extension):
    .pdf       PDF documents (pymupdf)           .ipynb     Jupyter notebooks (stdlib)
    .epub      E-books (ebooklib)                .pptx      PowerPoint (python-pptx)
    .docx      Word documents (python-docx)      .csv/.tsv  Tabular data (stdlib)
    .html      HTML files (stdlib)               .tex       LaTeX source (stdlib)
    .txt/.md   Plain text (stdlib)

    Optional: .png/.jpg (OCR via pytesseract), .mp3/.wav (transcription via whisper)

Dependencies: pip install pymupdf ebooklib python-docx python-pptx
              (stdlib: nbformat, html.parser, csv — no install needed)
"""
import asyncio
import csv
import html.parser as _html_parser
import io
import json
import mimetypes
import os
import re
import sys
import tempfile
from pathlib import Path
from urllib.parse import urlparse
from typing import Optional

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from sdk import DSSClient


# ── Format detection ────────────────────────────────────────────

EXTENSION_MAP = {
    '.pdf':  'pdf',   '.epub': 'epub',  '.docx': 'docx',
    '.pptx': 'pptx',  '.ipynb':'ipynb', '.html': 'html',
    '.htm':  'html',  '.txt':  'text',  '.md':   'markdown',
    '.tex':  'latex', '.csv':  'csv',   '.tsv':  'tsv',
    '.png':  'image', '.jpg':  'image', '.jpeg': 'image',
    '.mp3':  'audio', '.wav':  'audio', '.m4a':  'audio',
}

def detect_format(path: str | Path) -> str:
    suffix = Path(path).suffix.lower()
    return EXTENSION_MAP.get(suffix, 'unknown')


# ── Extractors ──────────────────────────────────────────────────

def extract_text(filepath: Path) -> str:
    return filepath.read_text(errors='replace')


def extract_markdown(filepath: Path) -> str:
    return filepath.read_text(errors='replace')


def extract_html(filepath: Path) -> str:
    """Extract text from HTML, stripping tags."""
    class TextExtractor(_html_parser.HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
            self.skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ('script', 'style', 'nav', 'footer', 'header'):
                self.skip = True
        def handle_endtag(self, tag):
            if tag in ('script', 'style', 'nav', 'footer', 'header'):
                self.skip = False
            if tag in ('p', 'br', 'div', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                self.text.append('\n')
        def handle_data(self, data):
            if not self.skip:
                self.text.append(data.strip())

    html_text = filepath.read_text(errors='replace')
    extractor = TextExtractor()
    extractor.feed(html_text)
    raw = '\n'.join(extractor.text)
    return _clean_whitespace(raw)


def extract_pdf(filepath: Path) -> str:
    try:
        import fitz  # pymupdf
    except ImportError:
        return _error("pymupdf not installed: pip install pymupdf")

    doc = fitz.open(str(filepath))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return '\n\n'.join(pages)


def extract_epub(filepath: Path) -> str:
    try:
        from ebooklib import epub
    except ImportError:
        return _error("ebooklib not installed: pip install ebooklib")

    book = epub.read_epub(str(filepath))
    chapters = []
    for item in book.get_items():
        if item.get_type() == 9:  # ITEM_DOCUMENT
            text = item.get_content().decode('utf-8', errors='replace')
            chapters.append(_strip_html_tags(text))
    return '\n\n'.join(chapters)


def extract_docx(filepath: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return _error("python-docx not installed: pip install python-docx")

    doc = Document(str(filepath))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return '\n\n'.join(paragraphs)


def extract_ipynb(filepath: Path) -> str:
    """Extract markdown + code cells from Jupyter notebook."""
    nb = json.loads(filepath.read_text())
    cells = []
    for cell in nb.get('cells', []):
        source = ''.join(cell.get('source', []))
        if cell.get('cell_type') == 'markdown':
            cells.append(source)
        elif cell.get('cell_type') == 'code':
            if source.strip():
                cells.append(f"```python\n{source}\n```")
    return '\n\n'.join(cells)


def extract_pptx(filepath: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return _error("python-pptx not installed: pip install python-pptx")

    prs = Presentation(str(filepath))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            slides.append(f"## Slide {i}\n\n" + '\n'.join(texts))
    return '\n\n---\n\n'.join(slides)


def extract_latex(filepath: Path) -> str:
    """Strip LaTeX commands, keep text."""
    text = filepath.read_text(errors='replace')
    # Remove comments
    text = re.sub(r'(?<!\\)%.*$', '', text, flags=re.MULTILINE)
    # Remove commands but keep their arguments' text content
    text = re.sub(r'\\[a-zA-Z]+\*?(?:\[[^\]]*\])*(?:\{[^}]*\})?', ' ', text)
    # Remove math mode
    text = re.sub(r'\$[^$]+\$', ' ', text)
    # Remove remaining braces
    text = text.replace('{', ' ').replace('}', ' ')
    return _clean_whitespace(text)


def extract_csv(filepath: Path, delimiter: str = ',') -> str:
    """Convert CSV to markdown table."""
    with open(filepath, newline='', errors='replace') as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = list(reader)
    if not rows:
        return ""
    # Limit to reasonable size
    rows = rows[:100]
    headers = rows[0]
    lines = ['| ' + ' | '.join(headers) + ' |']
    lines.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')
    for row in rows[1:]:
        padded = row + [''] * (len(headers) - len(row))
        lines.append('| ' + ' | '.join(padded[:len(headers)]) + ' |')
    return '\n'.join(lines)


def extract_tsv(filepath: Path) -> str:
    return extract_csv(filepath, delimiter='\t')


def extract_image_ocr(filepath: Path) -> str:
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        return _error("OCR deps not installed: pip install pillow pytesseract")

    img = Image.open(str(filepath))
    return pytesseract.image_to_string(img)


# ── Helpers ─────────────────────────────────────────────────────

def _strip_html_tags(html_text: str) -> str:
    class Stripper(_html_parser.HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
        def handle_data(self, data):
            self.text.append(data)
    s = Stripper()
    s.feed(html_text)
    return '\n'.join(s.text)


def _clean_whitespace(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    return text.strip()


def _error(msg: str) -> str:
    return f"[Error: {msg}]"


EXTRACTORS = {
    'text':     extract_text,
    'markdown': extract_markdown,
    'html':     extract_html,
    'pdf':      extract_pdf,
    'epub':     extract_epub,
    'docx':     extract_docx,
    'ipynb':    extract_ipynb,
    'pptx':     extract_pptx,
    'latex':    extract_latex,
    'csv':      extract_csv,
    'tsv':      extract_tsv,
    'image':    extract_image_ocr,
}


# ── Ingest ──────────────────────────────────────────────────────

async def ingest_file(dss: DSSClient, filepath: Path, fmt: str) -> bool:
    """Extract text and store in warehouse."""
    extractor = EXTRACTORS.get(fmt)
    if not extractor:
        print(f"  ⚠ Unsupported format: {fmt}")
        return False

    try:
        text = extractor(filepath)
    except Exception as e:
        print(f"  ✗ Extraction failed: {e}")
        return False

    if not text or text.startswith("[Error"):
        print(f"  ⚠ Empty or error content")
        return False

    title = filepath.stem[:200]
    word_count = len(text.split())

    payload = {
        "url": f"file://{filepath.absolute()}",
        "markdown": text[:100000],  # cap at 100KB
        "title": title,
        "source_domain": f"local:{fmt}",
        "word_count": word_count,
        "tags": ["local-file", fmt],
    }

    try:
        async with dss.client as client:
            resp = await client.post(
                f"{dss.warehouse}/ingest",
                json=payload,
                timeout=30.0,
            )
            resp.raise_for_status()
            print(f"  ✓ {fmt:6s} {title[:50]:50s} {word_count:>6d}w")
            return True
    except Exception as e:
        print(f"  ✗ Warehouse error: {e}")
        return False


async def process_path(dss: DSSClient, path: Path) -> int:
    """Process a file or directory, return count of ingested files."""
    if path.is_file():
        fmt = detect_format(path)
        ok = await ingest_file(dss, path, fmt)
        return 1 if ok else 0

    if path.is_dir():
        count = 0
        files = sorted(p for p in path.rglob('*') if p.is_file())
        for f in files:
            fmt = detect_format(f)
            if fmt == 'unknown':
                continue
            if await ingest_file(dss, f, fmt):
                count += 1
        return count

    print(f"  ⚠ Not found: {path}")
    return 0


async def process_url(dss: DSSClient, url: str) -> int:
    """Download remote file, extract, ingest."""
    import urllib.request
    parsed = urlparse(url)
    fname = Path(parsed.path).name or 'download'
    fmt = detect_format(fname)

    if fmt == 'unknown':
        print(f"  ⚠ Unknown format for URL: {url}")
        return 0

    with tempfile.NamedTemporaryFile(suffix=Path(fname).suffix, delete=False) as tmp:
        try:
            with urllib.request.urlopen(url, timeout=30) as resp:
                tmp.write(resp.read())
        except Exception as e:
            print(f"  ✗ Download failed: {e}")
            os.unlink(tmp.name)
            return 0

    ok = await ingest_file(dss, Path(tmp.name), fmt)
    os.unlink(tmp.name)
    return 1 if ok else 0


# ── Main ────────────────────────────────────────────────────────

async def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    paths = []
    urls = []
    args = sys.argv[1:]

    for i, arg in enumerate(args):
        if arg == '--dir' and i + 1 < len(args):
            paths.append(Path(args[i + 1]))
        elif arg == '--url' and i + 1 < len(args):
            urls.append(args[i + 1])
        elif not arg.startswith('--'):
            p = Path(arg)
            if p.exists():
                paths.append(p)
            elif arg.startswith('http'):
                urls.append(arg)

    if not paths and not urls:
        print("No files or URLs specified.")
        print(__doc__)
        sys.exit(1)

    async with DSSClient() as dss:
        total = 0

        for path in paths:
            print(f"\n{path}")
            count = await process_path(dss, path)
            total += count

        for url in urls:
            print(f"\n{url}")
            count = await process_url(dss, url)
            total += count

        print(f"\n✓ {total} files ingested")


if __name__ == "__main__":
    asyncio.run(main())
